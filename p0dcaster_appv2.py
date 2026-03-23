"""
PodcastLM Studio v2 — AI-powered podcast generator + study tools.
Bug fixes: thread-safe TTS, music ramp-up ordering, Whisper size check,
JSON mode fallback, smart truncation, fallback concat integrity.
Features: HD TTS, Edge TTS (free), per-speaker speed, script export,
session save/restore, SRT subtitles, DeepSeek as default LLM,
truncation recovery, NotebookLM-style study tools.
"""

import streamlit as st
import os
import re
import asyncio
import tempfile
import json
import requests
import io
import shutil
import hmac
import logging
from pathlib import Path
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import Optional, Tuple, List, Dict, Any

import PyPDF2
import docx
from pptx import Presentation
from bs4 import BeautifulSoup
import yt_dlp
import ffmpeg
import edge_tts
from openai import OpenAI

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
MAX_SOURCE_CHARS = 40_000
TTS_COST_PER_MILLION_CHARS = 15.0
TTS_HD_COST_PER_MILLION_CHARS = 30.0
AUDIO_BITRATE = "192k"
DEFAULT_BG_VOLUME = 0.12
MAX_TTS_WORKERS = 4
WHISPER_MAX_BYTES = 25 * 1024 * 1024

WORD_TARGETS = {
    "Short (2 min)": 800,
    "Medium (5 min)": 2_200,
    "Long (15 min)": 6_000,
    "Extra Long (30 min)": 12_000,
}

VOICE_MAP = {
    "Dynamic (Alloy & Nova)": ("alloy", "nova"),
    "Calm (Onyx & Shimmer)": ("onyx", "shimmer"),
    "Formal (Echo & Fable)": ("echo", "fable"),
}

EDGE_LANGUAGE_VOICES = {
    "English (US)": ("en-US-GuyNeural", "en-US-JennyNeural", "en-US-DavisNeural"),
    "English (UK)": ("en-GB-RyanNeural", "en-GB-SoniaNeural", "en-GB-ThomasNeural"),
    "Spanish": ("es-ES-AlvaroNeural", "es-ES-ElviraNeural", "es-MX-JorgeNeural"),
    "French": ("fr-FR-HenriNeural", "fr-FR-DeniseNeural", "fr-FR-YvesNeural"),
    "German": ("de-DE-ConradNeural", "de-DE-KatjaNeural", "de-DE-KillianNeural"),
    "Italian": ("it-IT-DiegoNeural", "it-IT-ElsaNeural", "it-IT-GiuseppeNeural"),
    "Portuguese": ("pt-BR-AntonioNeural", "pt-BR-FranciscaNeural", "pt-BR-ValerioNeural"),
    "Hindi": ("hi-IN-MadhurNeural", "hi-IN-SwaraNeural", "hi-IN-MadhurNeural"),
    "Urdu": ("ur-PK-AsadNeural", "ur-PK-UzmaNeural", "ur-PK-AsadNeural"),
    "Arabic": ("ar-SA-HamedNeural", "ar-SA-ZariyahNeural", "ar-EG-ShakirNeural"),
    "Hebrew": ("he-IL-AvriNeural", "he-IL-HilaNeural", "he-IL-AvriNeural"),
    "Russian": ("ru-RU-DmitryNeural", "ru-RU-SvetlanaNeural", "ru-RU-DmitryNeural"),
    "Turkish": ("tr-TR-AhmetNeural", "tr-TR-EmelNeural", "tr-TR-AhmetNeural"),
    "Japanese": ("ja-JP-KeitaNeural", "ja-JP-NanamiNeural", "ja-JP-KeitaNeural"),
    "Korean": ("ko-KR-InJoonNeural", "ko-KR-SunHiNeural", "ko-KR-InJoonNeural"),
    "Chinese (Mandarin)": ("zh-CN-YunxiNeural", "zh-CN-XiaoxiaoNeural", "zh-CN-YunjianNeural"),
    "Polish": ("pl-PL-MarekNeural", "pl-PL-ZofiaNeural", "pl-PL-MarekNeural"),
    "Dutch": ("nl-NL-MaartenNeural", "nl-NL-ColetteNeural", "nl-NL-MaartenNeural"),
    "Swedish": ("sv-SE-MattiasNeural", "sv-SE-SofieNeural", "sv-SE-MattiasNeural"),
    "Indonesian": ("id-ID-ArdiNeural", "id-ID-GadisNeural", "id-ID-ArdiNeural"),
    "Thai": ("th-TH-NiwatNeural", "th-TH-PremwadeeNeural", "th-TH-NiwatNeural"),
}

MUSIC_URLS = {
    "Lo-Fi (Study)": "https://cdn.pixabay.com/download/audio/2022/05/27/audio_1808fbf07a.mp3",
    "Upbeat (Morning)": "https://cdn.pixabay.com/download/audio/2024/05/24/audio_95e3f5f471.mp3",
    "Ambient (News)": "https://cdn.pixabay.com/download/audio/2022/03/10/audio_c8c8a73467.mp3",
    "Cinematic (Deep)": "https://cdn.pixabay.com/download/audio/2022/03/22/audio_c2b86c77ce.mp3",
}

SUPPORTED_LANGUAGES = [
    "English (US)", "English (UK)", "Spanish", "French", "German",
    "Italian", "Portuguese", "Hindi", "Urdu", "Arabic", "Hebrew",
    "Russian", "Turkish", "Japanese", "Korean", "Chinese (Mandarin)",
    "Polish", "Dutch", "Swedish", "Indonesian", "Thai",
]

NON_ENGLISH_LANGS = [
    "Urdu", "Arabic", "Hebrew", "Hindi", "Chinese",
    "Japanese", "Korean", "Russian", "Turkish",
]

DEEPSEEK_MODEL_MAP = {
    "DeepSeek-V3 (Recommended)": "deepseek-chat",
    "DeepSeek-R1 (Reasoning)": "deepseek-reasoner",
}

GROK_MODEL_MAP = {
    "Grok 4.1 Fast (Recommended)": "grok-4-1-fast-reasoning",
    "Grok 4 Full": "grok-4",
    "Grok 4 Fast": "grok-4-fast-reasoning",
    "Grok Code Fast": "grok-code-fast-1",
}

JSON_MODE_PREFIXES = ("gpt-", "grok-4-1", "deepseek-chat")

MODEL_MAX_TOKENS = {
    "deepseek-chat": 8_192,
    "deepseek-reasoner": 8_192,
    "gpt-4o-mini": 16_384,
}
DEFAULT_MODEL_MAX = 8_192

MAX_OUTPUT_TOKENS = {
    "Short (2 min)": 2_048,
    "Medium (5 min)": 4_096,
    "Long (15 min)": 8_192,
    "Extra Long (30 min)": 8_192,
}

SOURCE_CHARS_BY_LENGTH = {
    "Short (2 min)": 15_000,
    "Medium (5 min)": 30_000,
    "Long (15 min)": 40_000,
    "Extra Long (30 min)": 40_000,
}

# ---------------------------------------------------------------------------
# Study tool prompts
# ---------------------------------------------------------------------------
STUDY_TOOL_PROMPTS = {
    "Study Guide": (
        "Create a comprehensive study guide from the source material. Include:\n"
        "1. **Overview** — a 2–3 paragraph summary\n"
        "2. **Key Concepts** — the most important ideas, each with a brief explanation\n"
        "3. **Important Details** — specific facts, figures, dates, or findings worth remembering\n"
        "4. **Key Takeaways** — 5–10 bullet points summarizing what a reader should walk away knowing\n"
        "5. **Review Questions** — 5 questions a student could use to test their understanding\n\n"
        "Use clear Markdown formatting with headers, bold terms, and bullet points."
    ),
    "Briefing Document": (
        "Create an executive briefing document from the source material. Include:\n"
        "1. **Executive Summary** — 2–3 paragraphs covering the core message\n"
        "2. **Key Findings** — numbered list of the most important findings or arguments\n"
        "3. **Supporting Evidence** — specific data points, quotes, or references that back up the findings\n"
        "4. **Implications** — what these findings mean for the relevant field, organization, or audience\n"
        "5. **Recommendations** — actionable next steps based on the material\n"
        "6. **Open Questions** — areas that need further research or clarification\n\n"
        "Write in a professional, concise tone. Use Markdown formatting."
    ),
    "FAQ": (
        "Generate a comprehensive FAQ (Frequently Asked Questions) document from the source material.\n"
        "Create 12–15 question-and-answer pairs that cover:\n"
        "- The main topic and why it matters\n"
        "- Key concepts that might confuse a newcomer\n"
        "- Specific findings or claims and their evidence\n"
        "- Practical implications\n"
        "- Common misconceptions the source addresses\n\n"
        "Format each as:\n"
        "### Q: [question]\n"
        "**A:** [thorough 2–4 sentence answer]\n\n"
        "Use only information from the source. If something isn't addressed, don't invent answers."
    ),
    "Flashcards": (
        "Create a set of 20–30 flashcards from the source material for studying.\n"
        "Each flashcard should have a **Front** (term, question, or concept) and a "
        "**Back** (definition, answer, or explanation).\n\n"
        "Cover:\n"
        "- Key terms and definitions\n"
        "- Important facts and figures\n"
        "- Cause-and-effect relationships\n"
        "- Key people, places, or events mentioned\n\n"
        "Format as a JSON array:\n"
        '[{"front": "...", "back": "..."}, ...]\n\n'
        "Output ONLY the JSON array — no markdown fences, no extra text."
    ),
    "Timeline": (
        "Extract a chronological timeline of events, developments, or milestones "
        "from the source material.\n\n"
        "For each entry include:\n"
        "- **Date/Period** — as specific as the source allows\n"
        "- **Event** — what happened\n"
        "- **Significance** — why it matters (1 sentence)\n\n"
        "Format as a Markdown list:\n"
        "- **[date]** — [event]. *[significance]*\n\n"
        "If the source doesn't contain chronological information, create a logical "
        "progression of ideas/arguments instead, labeled as a 'Conceptual Progression.'\n"
        "Order from earliest to latest (or from foundational to advanced)."
    ),
    "Key Concepts": (
        "Create an alphabetized glossary of key concepts, terms, and definitions "
        "from the source material.\n\n"
        "For each entry include:\n"
        "- **Term** (bolded)\n"
        "- **Definition** — clear, concise explanation (1–3 sentences)\n"
        "- **Context** — how this term is used or why it matters in the source\n\n"
        "Format as:\n"
        "### A\n"
        "**Term** — Definition. *Context.*\n\n"
        "Include 15–30 terms. Only include terms that appear in or are directly "
        "relevant to the source material."
    ),
}

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Page config
# ---------------------------------------------------------------------------
st.set_page_config(
    page_title="PodcastLM Studio",
    page_icon="🎧",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ---------------------------------------------------------------------------
# Session state
# ---------------------------------------------------------------------------
if "session_start" not in st.session_state:
    st.session_state.session_start = datetime.now().strftime("%Y-%m-%d %H:%M")

_DEFAULTS: Dict[str, Any] = {
    "authenticated": False,
    "script_data": None,
    "source_text": "",
    "chat_history": [],
    "notebook_content": "",
    "rehearsal_audio": None,
    # Study tools
    "study_guide": None,
    "briefing_doc": None,
    "faq_doc": None,
    "flashcards": None,
    "timeline": None,
    "key_concepts": None,
}

for _k, _v in _DEFAULTS.items():
    if _k not in st.session_state:
        st.session_state[_k] = _v

if not st.session_state.notebook_content:
    st.session_state.notebook_content = (
        f"# Research Notebook\n"
        f"**Session Started:** {st.session_state.session_start}\n\n"
    )


# ---------------------------------------------------------------------------
# Authentication
# ---------------------------------------------------------------------------
def check_password() -> None:
    entered = st.session_state.get("password_input", "")
    expected = st.secrets.get("APP_PASSWORD", "")
    if hmac.compare_digest(entered, expected):
        st.session_state.authenticated = True
    else:
        st.error("Incorrect password.")


if not st.session_state.authenticated:
    st.title("🔐 Studio Login")
    st.text_input(
        "Enter Password", type="password",
        key="password_input", on_change=check_password,
    )
    st.stop()


# ---------------------------------------------------------------------------
# Utility functions
# ---------------------------------------------------------------------------
def smart_truncate(text: str, max_chars: int = MAX_SOURCE_CHARS) -> str:
    """Truncate at the last sentence boundary before max_chars."""
    if len(text) <= max_chars:
        return text
    truncated = text[:max_chars]
    for sep in [". ", ".\n", "! ", "!\n", "? ", "?\n"]:
        last = truncated.rfind(sep)
        if last > max_chars * 0.8:
            return truncated[: last + 1]
    return truncated


def _safe_line(line: Dict) -> Tuple[str, str]:
    """Extract speaker and text from a dialogue line, handling key variations."""
    speaker = (
        line.get("speaker")
        or line.get("Speaker")
        or line.get("SPEAKER")
        or "Unknown"
    )
    text = (
        line.get("text")
        or line.get("Text")
        or line.get("TEXT")
        or line.get("content")
        or line.get("line")
        or ""
    )
    return speaker, text


def repair_truncated_json(raw: str) -> Optional[Dict]:
    """Salvage usable script from truncated LLM JSON output."""
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        pass

    title_match = re.search(r'"title"\s*:\s*"((?:[^"\\]|\\.)*)"', raw)
    title = title_match.group(1) if title_match else "Untitled Podcast"

    pattern = r'\{\s*"speaker"\s*:\s*"([^"]*)"\s*,\s*"text"\s*:\s*"((?:[^"\\]|\\.)*)"\s*\}'
    matches = re.findall(pattern, raw, re.DOTALL)

    if not matches:
        return None

    dialogue = []
    for speaker, text in matches:
        try:
            unescaped = json.loads(f'"{text}"')
        except (json.JSONDecodeError, ValueError):
            unescaped = (
                text.replace('\\"', '"')
                .replace('\\n', '\n')
                .replace('\\t', '\t')
            )
        dialogue.append({"speaker": speaker, "text": unescaped})

    return {"title": title, "dialogue": dialogue}


def _seconds_to_srt_time(seconds: float) -> str:
    """Convert seconds to SRT timestamp HH:MM:SS,mmm."""
    hours = int(seconds // 3600)
    minutes = int((seconds % 3600) // 60)
    secs = int(seconds % 60)
    millis = int((seconds % 1) * 1000)
    return f"{hours:02d}:{minutes:02d}:{secs:02d},{millis:03d}"


def generate_srt(dialogue: List[Dict], words_per_minute: float = 150.0) -> str:
    """Generate SRT subtitle content with estimated timestamps."""
    srt_parts: List[str] = []
    current_time = 0.0

    for i, line in enumerate(dialogue):
        word_count = len(line["text"].split())
        duration = (word_count / words_per_minute) * 60.0

        start_str = _seconds_to_srt_time(current_time)
        end_str = _seconds_to_srt_time(current_time + duration)

        srt_parts.append(f"{i + 1}")
        srt_parts.append(f"{start_str} --> {end_str}")
        srt_parts.append(f"[{line['speaker']}] {line['text']}")
        srt_parts.append("")

        current_time += duration + 0.5

    return "\n".join(srt_parts)


def export_script_markdown(data: Dict) -> str:
    """Export script as readable Markdown."""
    lines = [f"# {data.get('title', 'Untitled Podcast')}\n"]
    for entry in data.get("dialogue", []):
        speaker, text = _safe_line(entry)
        lines.append(f"**{speaker}:** {text}\n")
    return "\n".join(lines)


def export_script_plain(data: Dict) -> str:
    """Export script as plain text."""
    lines = [data.get("title", "Untitled Podcast"), "=" * 40, ""]
    for entry in data.get("dialogue", []):
        speaker, text = _safe_line(entry)
        lines.append(f"[{speaker}] {text}")
        lines.append("")
    return "\n".join(lines)


def render_flashcards_markdown(cards: List[Dict]) -> str:
    """Convert flashcard dicts to readable Markdown."""
    lines = ["# Flashcards\n"]
    for i, card in enumerate(cards, 1):
        front = card.get("front") or card.get("Front") or card.get("question") or "?"
        back = card.get("back") or card.get("Back") or card.get("answer") or "?"
        lines.append(f"### Card {i}")
        lines.append(f"**Q:** {front}\n")
        lines.append(f"**A:** {back}\n")
        lines.append("---\n")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# LLM helpers
# ---------------------------------------------------------------------------
def get_llm_client(
    model_selection: str,
    specific_model_name: str,
    openai_key: str,
    xai_key: str,
    deepseek_key: str,
    budget_mode: bool,
) -> Tuple[Optional[OpenAI], Optional[str], Optional[str]]:
    """Return (client, model_name, error_message)."""
    if budget_mode or model_selection == "Model B (OpenAI)":
        if not openai_key:
            return None, None, "Missing OpenAI API Key"
        return OpenAI(api_key=openai_key), "gpt-4o-mini", None

    if model_selection == "Model A (DeepSeek) ⭐":
        if not deepseek_key:
            return None, None, "Missing DeepSeek API Key"
        actual = DEEPSEEK_MODEL_MAP.get(specific_model_name, "deepseek-chat")
        client = OpenAI(api_key=deepseek_key, base_url="https://api.deepseek.com")
        return client, actual, None

    if model_selection == "Model C (xAI Grok)":
        if not xai_key:
            return None, None, "Missing xAI API Key"
        actual = GROK_MODEL_MAP.get(specific_model_name, "grok-4-1-fast-reasoning")
        client = OpenAI(api_key=xai_key, base_url="https://api.x.ai/v1")
        return client, actual, None

    return None, None, "Invalid model selection"


def translate_if_needed(text: str, target_lang: str, openai_key: str) -> str:
    """Translate director notes via a dedicated OpenAI client."""
    if not any(lang in target_lang for lang in NON_ENGLISH_LANGS):
        return text
    if not openai_key:
        st.warning("OpenAI key required for translation — using original text.")
        return text
    try:
        oai = OpenAI(api_key=openai_key)
        res = oai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": f"Translate exactly to {target_lang}:\n\n{text}"}],
            max_tokens=1024,
        )
        return res.choices[0].message.content
    except Exception as e:
        st.warning(f"Translation failed, using original: {e}")
        return text


def generate_study_tool(
    client: OpenAI,
    model: str,
    source_text: str,
    tool_name: str,
    language: str,
    is_json: bool = False,
) -> Optional[str]:
    """Generate a study tool document from source material.
    Returns the raw text/markdown, or JSON string for flashcards."""
    prompt_body = STUDY_TOOL_PROMPTS.get(tool_name, "")
    if not prompt_body:
        return None

    prompt = (
        f"Based on the following source material, {prompt_body}\n\n"
        f"Write the output in {language}.\n\n"
        f"Source material:\n{smart_truncate(source_text)}"
    )

    desired_tokens = MODEL_MAX_TOKENS.get(model, DEFAULT_MODEL_MAX)

    kwargs: Dict[str, Any] = {
        "model": model,
        "messages": [{"role": "user", "content": prompt}],
        "max_tokens": desired_tokens,
    }

    if is_json and any(model.startswith(p) for p in JSON_MODE_PREFIXES):
        kwargs["response_format"] = {"type": "json_object"}

    try:
        res = client.chat.completions.create(**kwargs)
        raw = res.choices[0].message.content

        # Strip markdown fences if present
        cleaned = raw.strip()
        if cleaned.startswith("```"):
            cleaned = cleaned.split("\n", 1)[1]
        if cleaned.endswith("```"):
            cleaned = cleaned.rsplit("```", 1)[0]

        return cleaned.strip()
    except Exception as e:
        st.error(f"Generation failed: {e}")
        return None


# ---------------------------------------------------------------------------
# Audio helpers — Edge TTS
# ---------------------------------------------------------------------------
def _sync_edge_tts(text: str, voice: str, filepath: str, speed: float = 1.0) -> bool:
    """Synchronous wrapper for edge-tts."""
    try:
        rate_str = f"{int((speed - 1) * 100):+d}%"
        communicate = edge_tts.Communicate(text, voice, rate=rate_str)
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        try:
            loop.run_until_complete(communicate.save(filepath))
        finally:
            loop.close()
        p = Path(filepath)
        return p.exists() and p.stat().st_size > 0
    except Exception as e:
        logger.error("Edge TTS failed for voice %s: %s", voice, e)
        return False


# ---------------------------------------------------------------------------
# Audio helpers — OpenAI TTS
# ---------------------------------------------------------------------------
def generate_tts(
    client: Optional[OpenAI], text: str, voice: str, filepath: str,
    tts_model: str = "tts-1", speed: float = 1.0, use_edge: bool = False,
) -> bool:
    """Generate a single TTS segment (main-thread only)."""
    if use_edge:
        return _sync_edge_tts(text, voice, filepath, speed)
    try:
        response = client.audio.speech.create(
            model=tts_model, voice=voice, input=text, speed=speed,
        )
        response.stream_to_file(filepath)
        p = Path(filepath)
        return p.exists() and p.stat().st_size > 0
    except Exception as e:
        logger.error("TTS failed for voice %s: %s", voice, e)
        st.warning(f"TTS failed for voice '{voice}': {e}")
        return False


def _voice_line_worker(args: tuple) -> Tuple[int, bool, Optional[str]]:
    """Thread-safe TTS worker — NO Streamlit calls allowed here."""
    (i, line, tts_client, m_voice, f_voice, c_voice, tmp_path,
     tts_model_name, h1_speed, h2_speed, c_speed, use_edge) = args

    voice = m_voice if line["speaker"] == "Host 1" else f_voice
    speed = h1_speed if line["speaker"] == "Host 1" else h2_speed
    if line["speaker"] == "Caller":
        voice = c_voice
        speed = c_speed

    path = str(tmp_path / f"{i}.mp3")

    try:
        if use_edge:
            ok = _sync_edge_tts(line["text"], voice, path, speed)
        else:
            response = tts_client.audio.speech.create(
                model=tts_model_name, voice=voice, input=line["text"], speed=speed,
            )
            response.stream_to_file(path)
            p = Path(path)
            ok = p.exists() and p.stat().st_size > 0

        if ok:
            return i, True, None
        return i, False, f"Empty audio file for line {i}"
    except Exception as e:
        return i, False, str(e)


# ---------------------------------------------------------------------------
# Audio helpers — effects and mixing
# ---------------------------------------------------------------------------
def apply_phone_effect(input_path: str, output_path: str) -> None:
    """Bandpass + echo to simulate a phone caller."""
    try:
        main = (
            ffmpeg.input(input_path)
            .filter("lowpass", f=3000)
            .filter("highpass", f=300)
        )
        echo = (
            ffmpeg.input(input_path)
            .filter("adelay", delays="120|120")
            .filter("volume", "0.7")
        )
        mixed = ffmpeg.filter([main, echo], "amix", inputs=2)
        ffmpeg.output(mixed, output_path, acodec="mp3", audio_bitrate=AUDIO_BITRATE).run(
            overwrite_output=True, quiet=True,
        )
    except Exception as e:
        logger.warning("Phone effect failed, copying original: %s", e)
        shutil.copy(input_path, output_path)


def download_file(url: str, save_path: str) -> bool:
    """Download a file with a browser-like UA header."""
    try:
        r = requests.get(
            url, headers={"User-Agent": "Mozilla/5.0"}, stream=True, timeout=30,
        )
        r.raise_for_status()
        with open(save_path, "wb") as f:
            for chunk in r.iter_content(8192):
                f.write(chunk)
        return True
    except Exception as e:
        st.warning(f"Download failed: {e}")
        return False


def _write_uploaded_to_disk(uploaded_file, dest: str) -> None:
    """Persist a Streamlit UploadedFile to a real path for ffmpeg."""
    with open(dest, "wb") as f:
        f.write(uploaded_file.getvalue())


def mix_final_audio(
    tmp_dir: str,
    script_dialogue: List[Dict],
    bg_source: str,
    selected_bg_url: Optional[str],
    uploaded_bg_file,
    music_ramp_up: bool,
    uploaded_intro,
    uploaded_outro,
) -> Optional[Path]:
    """Combine voiced segments, background music, intro/outro into one MP3."""
    tmp = Path(tmp_dir)
    inputs = []

    for i, line in enumerate(script_dialogue):
        seg_path = tmp / f"{i}.mp3"
        if not seg_path.exists() or seg_path.stat().st_size == 0:
            st.warning(f"Skipping missing/empty segment {i}.")
            continue
        if line["speaker"] == "Caller":
            phone_path = tmp / f"phone_{i}.mp3"
            apply_phone_effect(str(seg_path), str(phone_path))
            inputs.append(ffmpeg.input(str(phone_path)))
        else:
            inputs.append(ffmpeg.input(str(seg_path)))

    if not inputs:
        st.error("No valid audio segments to mix.")
        return None

    if len(inputs) > 1:
        dialogue = ffmpeg.concat(*inputs, v=0, a=1, n=len(inputs))
    else:
        dialogue = inputs[0]

    dialogue = dialogue.filter("loudnorm", I=-16, LRA=11, TP=-1.5)

    if music_ramp_up and bg_source != "None":
        silence = ffmpeg.input(
            "anullsrc=channel_layout=stereo:sample_rate=44100", f="lavfi", t=5,
        )
        dialogue = ffmpeg.concat(silence, dialogue, v=0, a=1)

    if bg_source != "None":
        bg_path = tmp / "bg.mp3"
        if bg_source == "Presets" and selected_bg_url:
            download_file(selected_bg_url, str(bg_path))
        elif uploaded_bg_file:
            _write_uploaded_to_disk(uploaded_bg_file, str(bg_path))

        if bg_path.exists() and bg_path.stat().st_size > 0:
            bg = ffmpeg.input(str(bg_path))
            bg = bg.filter("aloop", loop=-1, size="2147483647")
            bg = bg.filter("volume", DEFAULT_BG_VOLUME)
            dialogue = ffmpeg.filter(
                [bg, dialogue], "amix", inputs=2, duration="shortest",
            )

    if uploaded_intro:
        intro_path = str(tmp / "intro.mp3")
        _write_uploaded_to_disk(uploaded_intro, intro_path)
        dialogue = ffmpeg.concat(ffmpeg.input(intro_path), dialogue, v=0, a=1)

    if uploaded_outro:
        outro_path = str(tmp / "outro.mp3")
        _write_uploaded_to_disk(uploaded_outro, outro_path)
        dialogue = ffmpeg.concat(dialogue, ffmpeg.input(outro_path), v=0, a=1)

    out_path = tmp / "podcast.mp3"
    try:
        ffmpeg.output(dialogue, str(out_path), acodec="mp3", audio_bitrate=AUDIO_BITRATE).run(
            overwrite_output=True, quiet=True,
        )
    except ffmpeg.Error as e:
        stderr = e.stderr.decode(errors="ignore") if e.stderr else "unknown"
        st.warning(f"Advanced mix failed — falling back.\n```\n{stderr}\n```")
        try:
            simple = ffmpeg.concat(*inputs, v=0, a=1, n=len(inputs))
            ffmpeg.output(simple, str(out_path), acodec="mp3", audio_bitrate=AUDIO_BITRATE).run(
                overwrite_output=True, quiet=True,
            )
        except Exception as e2:
            st.error(f"Fallback concat also failed: {e2}")
            return None

    return out_path if out_path.exists() else None


# ---------------------------------------------------------------------------
# Content extraction
# ---------------------------------------------------------------------------
@st.cache_data(show_spinner=False)
def scrape_website(url: str) -> Optional[str]:
    """Fetch and extract body text from a web page."""
    try:
        r = requests.get(url, headers={"User-Agent": "Mozilla/5.0"}, timeout=15)
        r.raise_for_status()
        soup = BeautifulSoup(r.content, "html.parser")
        for tag in soup(["script", "style", "header", "footer", "nav", "aside"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        return text if text else None
    except Exception as e:
        st.warning(f"Scraping failed: {e}")
        return None


def extract_text_from_files(files, audio_client: Optional[OpenAI] = None) -> str:
    """Extract text from PDFs, DOCX, PPTX, TXT, and audio files."""
    parts: List[str] = []
    for file in files:
        try:
            name = file.name.lower()
            raw = file.getvalue()

            if name.endswith(".pdf"):
                reader = PyPDF2.PdfReader(io.BytesIO(raw))
                parts.extend(page.extract_text() or "" for page in reader.pages)

            elif name.endswith(".docx"):
                doc = docx.Document(io.BytesIO(raw))
                parts.extend(p.text for p in doc.paragraphs)

            elif name.endswith(".pptx"):
                prs = Presentation(io.BytesIO(raw))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            parts.append(shape.text)

            elif name.endswith(".txt"):
                parts.append(raw.decode("utf-8", errors="replace"))

            elif name.endswith((".mp3", ".wav", ".m4a", ".mp4", ".webm")):
                if not audio_client:
                    st.warning(f"OpenAI key required to transcribe {file.name}")
                elif len(raw) > WHISPER_MAX_BYTES:
                    st.warning(
                        f"⚠️ {file.name} is {len(raw) / 1_048_576:.1f} MB — "
                        f"exceeds Whisper's 25 MB limit. Skipping."
                    )
                else:
                    with st.spinner(f"Transcribing {file.name}…"):
                        transcript = audio_client.audio.transcriptions.create(
                            model="whisper-1", file=(file.name, raw),
                        )
                        parts.append(transcript.text)
            else:
                st.warning(f"Unsupported file type: {file.name}")
        except Exception as e:
            st.error(f"Error reading {file.name}: {e}")

    return "\n".join(parts)


def download_and_transcribe_video(
    url: str, audio_client: OpenAI,
) -> Tuple[Optional[str], Optional[str]]:
    """Download audio from a video URL and transcribe via Whisper."""
    try:
        with tempfile.TemporaryDirectory() as tmp_dir:
            ydl_opts = {
                "format": "bestaudio/best",
                "outtmpl": os.path.join(tmp_dir, "audio.%(ext)s"),
                "postprocessors": [
                    {"key": "FFmpegExtractAudio", "preferredcodec": "mp3"},
                ],
                "quiet": True,
                "http_headers": {"User-Agent": "Mozilla/5.0"},
            }
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                ydl.download([url])

            audio_path = next(Path(tmp_dir).glob("audio.*"))
            with open(audio_path, "rb") as f:
                transcript = audio_client.audio.transcriptions.create(
                    model="whisper-1", file=f,
                )
            return transcript.text, None
    except Exception as e:
        return None, str(e)


# ---------------------------------------------------------------------------
# Sidebar
# ---------------------------------------------------------------------------
with st.sidebar:
    st.title("⚙️ Studio Settings")

    openai_key = st.secrets.get("OPENAI_API_KEY") or st.text_input(
        "OpenAI API Key", type="password",
    )
    deepseek_key = st.secrets.get("DEEPSEEK_API_KEY") or st.text_input(
        "DeepSeek API Key", type="password",
    )
    xai_key = st.secrets.get("XAI_API_KEY") or st.text_input(
        "xAI API Key (Optional)", type="password",
    )

    model_choice = st.radio(
        "Intelligence Engine",
        ["Model A (DeepSeek) ⭐", "Model B (OpenAI)", "Model C (xAI Grok)"],
        help="DeepSeek-V3 is the default — best quality-to-cost ratio",
    )
    deepseek_version = "DeepSeek-V3 (Recommended)"
    xai_version = "Grok 4.1 Fast (Recommended)"
    if model_choice == "Model A (DeepSeek) ⭐":
        deepseek_version = st.selectbox("DeepSeek Model", list(DEEPSEEK_MODEL_MAP.keys()))
    elif model_choice == "Model C (xAI Grok)":
        xai_version = st.selectbox("Grok Model", list(GROK_MODEL_MAP.keys()))

    budget_mode = st.checkbox(
        "💰 Budget Mode (GPT-4o-mini)",
        help="Overrides model selection — uses cheapest option",
    )
    privacy_mode = st.toggle("🔒 Privacy Mode", value=False)

    if st.button("🔄 New Session"):
        for key in _DEFAULTS:
            st.session_state[key] = _DEFAULTS[key]
        st.session_state.session_start = datetime.now().strftime("%Y-%m-%d %H:%M")
        st.session_state.notebook_content = (
            f"# Research Notebook\n**Session Started:** {st.session_state.session_start}\n\n"
        )
        st.rerun()

    st.divider()
    st.subheader("🌐 Language & Length")
    language = st.selectbox("Output Language", SUPPORTED_LANGUAGES)
    length_option = st.select_slider(
        "Duration", list(WORD_TARGETS.keys()), value="Medium (5 min)",
    )

    # --- TTS Provider & Voice Selection ---
    st.divider()
    st.subheader("🎙️ Voices")
    tts_provider = st.radio(
        "TTS Provider",
        ["OpenAI", "Edge TTS (Free)"],
        horizontal=True,
        help="Edge TTS uses Microsoft's neural voices — completely free, no API key needed",
    )

    use_edge = tts_provider == "Edge TTS (Free)"

    if not use_edge:
        voice_style = st.selectbox("Voice Pair", list(VOICE_MAP.keys()))
        m_voice, f_voice = VOICE_MAP[voice_style]
        c_voice = "fable"

        tts_hd = st.checkbox(
            "🔊 HD Voices (TTS-1-HD)",
            help="2× TTS cost, noticeably better quality",
        )
        tts_model = "tts-1-hd" if tts_hd else "tts-1"
    else:
        edge_voices = EDGE_LANGUAGE_VOICES.get(
            language, EDGE_LANGUAGE_VOICES["English (US)"],
        )
        m_voice, f_voice, c_voice = edge_voices
        tts_model = "edge"
        tts_hd = False

        st.success("🆓 Free — no TTS charges")
        st.caption(f"Voices auto-matched for **{language}**:")
        st.caption(f"🗣️ Host 1: `{m_voice}`")
        st.caption(f"🗣️ Host 2: `{f_voice}`")
        if c_voice != m_voice:
            st.caption(f"📞 Caller: `{c_voice}`")

    st.caption("Speaking Speed")
    spd_c1, spd_c2 = st.columns(2)
    with spd_c1:
        host1_speed = st.slider("Host 1", 0.70, 1.30, 1.0, 0.05, key="h1_speed")
    with spd_c2:
        host2_speed = st.slider("Host 2", 0.70, 1.30, 1.0, 0.05, key="h2_speed")
    caller_speed = st.slider("Caller", 0.70, 1.30, 0.95, 0.05, key="caller_speed")

    st.subheader("🎙️ Hosts")
    host1_persona = st.text_input("Host 1 Persona", "Male, curious, slightly skeptical")
    host2_persona = st.text_input("Host 2 Persona", "Female, enthusiastic expert")

    st.divider()
    st.subheader("🎵 Music")
    bg_source = st.radio("Background Music", ["Presets", "Upload Custom", "None"], horizontal=True)
    music_ramp_up = st.checkbox("Start music 5 s early")
    selected_bg_url = None
    uploaded_bg_file = None
    if bg_source == "Presets":
        music_choice = st.selectbox("Track", list(MUSIC_URLS.keys()))
        selected_bg_url = MUSIC_URLS[music_choice]
    elif bg_source == "Upload Custom":
        uploaded_bg_file = st.file_uploader("Upload Loop", type=["mp3", "wav"])

    with st.expander("Intro / Outro"):
        uploaded_intro = st.file_uploader("Intro clip", type=["mp3", "wav"])
        uploaded_outro = st.file_uploader("Outro clip", type=["mp3", "wav"])

    # Session save / restore
    st.divider()
    with st.expander("💾 Session Save / Restore"):
        if st.session_state.script_data or st.session_state.source_text:
            session_export = {
                "script_data": st.session_state.script_data,
                "source_text": st.session_state.source_text,
                "chat_history": st.session_state.chat_history,
                "study_guide": st.session_state.study_guide,
                "briefing_doc": st.session_state.briefing_doc,
                "faq_doc": st.session_state.faq_doc,
                "flashcards": st.session_state.flashcards,
                "timeline": st.session_state.timeline,
                "key_concepts": st.session_state.key_concepts,
                "exported_at": datetime.now().isoformat(),
            }
            st.download_button(
                "⬇️ Download Session",
                json.dumps(session_export, indent=2, ensure_ascii=False),
                file_name=f"session_{datetime.now().strftime('%Y%m%d_%H%M')}.json",
                mime="application/json",
            )
        else:
            st.caption("Nothing to save yet.")

        uploaded_session = st.file_uploader(
            "Upload saved session", type=["json"], key="session_upload",
        )
        if uploaded_session:
            try:
                restored = json.loads(uploaded_session.getvalue().decode("utf-8"))
                if st.button("✅ Restore Session"):
                    st.session_state.script_data = restored.get("script_data")
                    st.session_state.source_text = restored.get("source_text", "")
                    st.session_state.chat_history = restored.get("chat_history", [])
                    st.session_state.study_guide = restored.get("study_guide")
                    st.session_state.briefing_doc = restored.get("briefing_doc")
                    st.session_state.faq_doc = restored.get("faq_doc")
                    st.session_state.flashcards = restored.get("flashcards")
                    st.session_state.timeline = restored.get("timeline")
                    st.session_state.key_concepts = restored.get("key_concepts")
                    st.success(
                        f"Session restored (saved {restored.get('exported_at', 'unknown')})"
                    )
                    st.rerun()
            except Exception as e:
                st.error(f"Invalid session file: {e}")

    # --- Cost estimate ---
    st.divider()
    st.subheader("💵 Cost Estimate")

    if use_edge:
        tts_rate = 0.0
    else:
        tts_rate = TTS_HD_COST_PER_MILLION_CHARS if tts_hd else TTS_COST_PER_MILLION_CHARS

    def _get_llm_cost() -> float:
        if budget_mode or model_choice == "Model B (OpenAI)":
            return 0.10
        elif model_choice == "Model A (DeepSeek) ⭐":
            return 0.03
        else:
            return 0.30

    if st.session_state.script_data:
        total_chars = sum(
            len(l["text"]) for l in st.session_state.script_data["dialogue"]
        )
        total_lines = len(st.session_state.script_data["dialogue"])
        tts_cost = (total_chars / 1_000_000) * tts_rate
        llm_cost = _get_llm_cost()
        c1, c2 = st.columns(2)
        c1.metric("TTS", f"${tts_cost:.3f}" if not use_edge else "FREE")
        c2.metric("LLM", f"${llm_cost:.2f}")
        st.success(f"**Estimated total ≈ ${tts_cost + llm_cost:.2f}**")
        st.caption(
            f"{total_lines} lines · {total_chars:,} chars"
            + (" · HD" if tts_hd else "")
            + (" · Edge TTS" if use_edge else "")
        )
    else:
        target_words = WORD_TARGETS[length_option]
        est_chars = int(target_words * 5.5)
        est_tts = (est_chars / 1_000_000) * tts_rate
        est_llm = _get_llm_cost()
        label = f"Pre-gen estimate ≈ **${est_tts + est_llm:.2f}** for {length_option}"
        if use_edge:
            label += " · **TTS is free**"
        elif tts_hd:
            label += " · HD"
        st.info(label)


# ---------------------------------------------------------------------------
# Helper to resolve the correct specific_model_name
# ---------------------------------------------------------------------------
def _resolve_model_name() -> str:
    if model_choice == "Model A (DeepSeek) ⭐":
        return deepseek_version
    elif model_choice == "Model C (xAI Grok)":
        return xai_version
    else:
        return ""


# ---------------------------------------------------------------------------
# Shared OpenAI client (for Whisper transcription — always OpenAI)
# ---------------------------------------------------------------------------
audio_client: Optional[OpenAI] = OpenAI(api_key=openai_key) if openai_key else None

# ---------------------------------------------------------------------------
# Main tabs
# ---------------------------------------------------------------------------
st.title("🎧 PodcastLM Studio")
# ---------------------------------------------------------------------------
# In-app guide
# ---------------------------------------------------------------------------
with st.expander("📘 How to Use This App", expanded=False):
    st.markdown("""
### Welcome to PodcastLM Studio!

This app turns any document, article, video, or text into a fully produced AI podcast — with study tools, research chat, and more. Here's how each tab works:

---

#### 📄 Tab 1: Source
Upload your content. This is the foundation for everything else in the app.

- **Files** — PDF, DOCX, PPTX, TXT, or audio/video files (auto-transcribed via Whisper)
- **Web URL** — paste an article link and the app scrapes the text
- **Video URL** — paste a YouTube link and the app downloads + transcribes the audio
- **Text** — paste raw text directly

> **Tip:** You can upload multiple files at once. The app combines them into a single source.

---

#### 💬 Tab 2: Research Chat
Ask questions about your source material. The AI answers using **only** what you uploaded — no hallucination, no outside knowledge.

Great for:
- Quickly understanding a long document
- Finding specific facts or figures
- Exploring the source before generating a podcast

---

#### 📝 Tab 3: Script & Rehearsal
Generate a two-host podcast script from your source.

1. **Director Notes** — tell the AI what tone, focus, or style you want (e.g., "Make it funny," "Focus on the methodology")
2. **Caller Question** — optionally add a phone-in caller who asks a specific question
3. Click **✨ Generate Script** and wait 30–90 seconds
4. **Edit** any line, change speakers, rewrite dialogue
5. **Rehearse** — preview any single line with the selected voice before producing

> **Tip:** Use the sidebar to change language, duration, voice pair, and speaking speed.

---

#### 🎚️ Tab 4: Produce
Turn your script into a finished podcast MP3.

- All dialogue lines are voiced in parallel (fast!)
- Background music, intro/outro, and caller phone effects are mixed in automatically
- Download the final MP3 and SRT subtitles when done

> **Tip:** Switch to **Edge TTS (Free)** in the sidebar to eliminate all voice generation costs.

---

#### 📚 Tab 5: Study Tools
Generate NotebookLM-style study materials from your source:

| Tool | What it creates |
|---|---|
| **Study Guide** | Structured summary with key concepts, takeaways, and review questions |
| **Briefing Document** | Executive-style brief with findings, implications, and recommendations |
| **FAQ** | 12–15 Q&A pairs covering the most important topics |
| **Flashcards** | Interactive front/back cards for memorization |
| **Timeline** | Chronological events or conceptual progression |
| **Key Concepts** | Alphabetized glossary of important terms |

You can generate them one at a time or click **🚀 Generate All** to create everything at once. All tools are downloadable as Markdown.

---

#### ⚙️ Sidebar Settings

| Setting | What it does |
|---|---|
| **Intelligence Engine** | Choose your LLM: DeepSeek (cheapest), OpenAI, or Grok |
| **Budget Mode** | Forces GPT-4o-mini for maximum savings |
| **Privacy Mode** | Wipes source text from memory after script generation |
| **TTS Provider** | OpenAI (paid, 6 voices) or Edge TTS (free, 300+ voices) |
| **Voice Pair** | Choose the host voice combination |
| **Speaking Speed** | Per-speaker speed sliders (Host 1, Host 2, Caller) |
| **HD Voices** | Higher quality OpenAI TTS (2× cost) |
| **Background Music** | Preset tracks, custom upload, or none |
| **Session Save/Restore** | Download your session as JSON; upload it later to continue |

---

#### 💵 Cost

| Component | Billed to | Typical cost |
|---|---|---|
| Script generation | DeepSeek / OpenAI / Grok key | ~$0.03–$0.30 |
| Voice generation | OpenAI key (or free with Edge TTS) | $0.00–$1.00 |
| Study tools | Same as script generation | ~$0.01–$0.05 each |

**Cheapest setup:** DeepSeek + Edge TTS = **~$0.03 per podcast**.

---

*Need help? Open an issue on the [GitHub repo](https://github.com/your-username/podcastlm-studio).*
    """)
tab1, tab2, tab3, tab4, tab5 = st.tabs(
    ["📄 Source", "💬 Research Chat", "📝 Script & Rehearsal", "🎚️ Produce", "📚 Study Tools"],
)


# === TAB 1 — SOURCE =========================================================
with tab1:
    st.info("Upload content — this drives the research chat, podcast, and study tools.")
    input_type = st.radio(
        "Input Type", ["Files", "Web URL", "Video URL", "Text"], horizontal=True,
    )
    new_text = ""

    if input_type == "Files":
        files = st.file_uploader(
            "Upload documents or audio",
            accept_multiple_files=True,
            type=["pdf", "docx", "pptx", "txt", "mp3", "wav", "m4a", "mp4", "webm"],
        )
        if files and st.button("Process Files"):
            with st.spinner("Extracting text…"):
                new_text = extract_text_from_files(files, audio_client)

    elif input_type == "Web URL":
        url = st.text_input("Article URL")
        if url and st.button("Scrape"):
            with st.spinner("Scraping…"):
                new_text = scrape_website(url) or ""

    elif input_type == "Video URL":
        vid_url = st.text_input("YouTube / Video URL")
        if vid_url and st.button("Transcribe"):
            if not audio_client:
                st.error("OpenAI key required for transcription.")
            else:
                with st.spinner("Downloading & transcribing video…"):
                    text, err = download_and_transcribe_video(vid_url, audio_client)
                    if err:
                        st.error(f"Transcription error: {err}")
                    new_text = text or ""

    elif input_type == "Text":
        new_text = st.text_area("Paste text", height=300)

    if new_text and new_text != st.session_state.source_text:
        st.session_state.source_text = new_text
        st.session_state.chat_history = []
        st.session_state.notebook_content += (
            f"\n---\n### New Source ({datetime.now().strftime('%H:%M')})\n\n"
        )
        # Clear study tools when source changes
        for tool_key in ["study_guide", "briefing_doc", "faq_doc", "flashcards", "timeline", "key_concepts"]:
            st.session_state[tool_key] = None
        st.success(f"✅ Source loaded — {len(new_text):,} characters.")

    if st.session_state.source_text:
        with st.expander(
            f"📖 Current source ({len(st.session_state.source_text):,} chars)",
            expanded=False,
        ):
            preview = st.session_state.source_text[:2000]
            if len(st.session_state.source_text) > 2000:
                preview += "…"
            st.text(preview)


# === TAB 2 — RESEARCH CHAT ==================================================
with tab2:
    st.header("💬 Research Chat")
    if not st.session_state.source_text:
        st.info("Load source content in the **Source** tab first.")
    else:
        for entry in st.session_state.chat_history:
            role = "user" if entry["role"] == "user" else "assistant"
            with st.chat_message(role):
                st.markdown(entry["content"])

        user_question = st.chat_input("Ask about the source material…")
        if user_question:
            st.session_state.chat_history.append(
                {"role": "user", "content": user_question},
            )
            with st.chat_message("user"):
                st.markdown(user_question)

            client, model, err = get_llm_client(
                model_choice, _resolve_model_name(),
                openai_key, xai_key, deepseek_key, budget_mode,
            )
            if err:
                st.error(err)
            else:
                messages = [
                    {
                        "role": "system",
                        "content": (
                            "You are a helpful research assistant. Answer questions using "
                            "ONLY the provided source material. If the answer isn't in the "
                            "source, say so."
                        ),
                    },
                    {
                        "role": "system",
                        "content": f"Source:\n{smart_truncate(st.session_state.source_text)}",
                    },
                    *st.session_state.chat_history,
                ]
                with st.chat_message("assistant"):
                    with st.spinner("Thinking…"):
                        try:
                            response = client.chat.completions.create(
                                model=model, messages=messages, max_tokens=1024,
                            )
                            ai_reply = response.choices[0].message.content
                            st.markdown(ai_reply)
                            st.session_state.chat_history.append(
                                {"role": "assistant", "content": ai_reply},
                            )
                        except Exception as e:
                            st.error(f"LLM error: {e}")

        if st.button("🗑️ Clear Chat"):
            st.session_state.chat_history = []
            st.rerun()


# === TAB 3 — SCRIPT GENERATION & REHEARSAL ===================================
with tab3:
    col_dir, col_call = st.columns(2)
    with col_dir:
        user_instructions = st.text_area(
            "🎬 Director Notes",
            placeholder="e.g., Make it funny, focus on the key findings",
        )
    with col_call:
        caller_prompt = st.text_area(
            "📞 Caller Question (optional)",
            placeholder="e.g., What does this mean for everyday people?",
        )

    if st.button("✨ Generate Script", type="primary"):
        if not st.session_state.source_text:
            st.error("Load source content first (Tab 1).")
        else:
            client, model, err = get_llm_client(
                model_choice, _resolve_model_name(),
                openai_key, xai_key, deepseek_key, budget_mode,
            )
            if err:
                st.error(err)
            else:
                with st.spinner("Writing script… this may take 30–90 seconds."):
                    target_words = WORD_TARGETS[length_option]
                    translated = translate_if_needed(
                        user_instructions, language, openai_key,
                    )

                    call_in = ""
                    if caller_prompt:
                        call_in = (
                            f'\nInclude a "Caller" speaker who asks: \'{caller_prompt}\' '
                            f"— the hosts then respond thoughtfully."
                        )

                    source_limit = SOURCE_CHARS_BY_LENGTH.get(length_option, MAX_SOURCE_CHARS)
                    source_text = smart_truncate(
                        st.session_state.source_text, max_chars=source_limit,
                    )

                    prompt = f"""Create a podcast script in {language}.

Host 1 persona: {host1_persona}
Host 2 persona: {host2_persona}

Write a very detailed, natural, conversational podcast script with approximately \
{target_words} total words ({length_option}).
Use long explanations, tangents, humor, and back-and-forth dialogue. \
NEVER truncate or summarize lines.

Director notes: {translated}
{call_in}

Output **strict JSON only** — no markdown fences:
{{"title": "...", "dialogue": [{{"speaker": "Host 1", "text": "..."}}, \
{{"speaker": "Host 2", "text": "..."}}, ...]}}

Source material:
{source_text}"""

                    try:
                        desired_tokens = MAX_OUTPUT_TOKENS.get(length_option, 4_096)
                        model_cap = MODEL_MAX_TOKENS.get(model, DEFAULT_MODEL_MAX)
                        safe_tokens = min(desired_tokens, model_cap)

                        kwargs: Dict[str, Any] = {
                            "model": model,
                            "messages": [{"role": "user", "content": prompt}],
                            "max_tokens": safe_tokens,
                        }
                        if any(model.startswith(p) for p in JSON_MODE_PREFIXES):
                            kwargs["response_format"] = {"type": "json_object"}

                        res = client.chat.completions.create(**kwargs)
                        raw = res.choices[0].message.content

                        finish_reason = getattr(res.choices[0], "finish_reason", "stop")
                        was_truncated = finish_reason in ("length", "max_tokens")

                        cleaned = raw.strip()
                        if cleaned.startswith("```"):
                            cleaned = cleaned.split("\n", 1)[1]
                        if cleaned.endswith("```"):
                            cleaned = cleaned.rsplit("```", 1)[0]
                        cleaned = cleaned.strip()

                        parsed = None
                        try:
                            parsed = json.loads(cleaned)
                        except json.JSONDecodeError:
                            parsed = repair_truncated_json(cleaned)

                        if parsed and "dialogue" in parsed:
                            normalized = []
                            for entry in parsed["dialogue"]:
                                speaker, text = _safe_line(entry)
                                if text:
                                    normalized.append({"speaker": speaker, "text": text})
                            parsed["dialogue"] = normalized

                        if parsed and "dialogue" in parsed and len(parsed["dialogue"]) > 0:
                            st.session_state.script_data = parsed

                            if was_truncated:
                                actual_words = sum(
                                    len(l["text"].split()) for l in parsed["dialogue"]
                                )
                                st.warning(
                                    f"⚠️ Output was truncated by the model — "
                                    f"recovered {len(parsed['dialogue'])} lines "
                                    f"(~{actual_words:,} words out of "
                                    f"{target_words:,} requested).\n\n"
                                    f"**To fix:** try a shorter Duration, or switch to "
                                    f"Budget Mode (GPT-4o-mini supports longer outputs)."
                                )
                            else:
                                st.success("✅ Script generated!")

                            if privacy_mode:
                                st.session_state.source_text = ""
                                st.info("🔒 Source text wiped (Privacy Mode).")
                        else:
                            st.error(
                                f"Could not parse or repair the script output.\n\n"
                                f"Raw output (first 500 chars):\n"
                                f"```\n{raw[:500]}\n```\n\n"
                                f"**Try:** shorter Duration, or Budget Mode (GPT-4o-mini)."
                            )

                    except Exception as e:
                        st.error(f"Script generation failed: {e}")

    # --- Display, edit, export script ---
    if st.session_state.script_data:
        data = st.session_state.script_data
        dialogue = data.get("dialogue", [])
        word_count = sum(len(l["text"].split()) for l in dialogue)

        st.subheader(data.get("title", "Untitled Podcast"))
        st.caption(
            f"{len(dialogue)} lines · ~{word_count:,} words · "
            f"est. {word_count // 150} min"
        )

        exp_c1, exp_c2, exp_c3 = st.columns(3)
        with exp_c1:
            st.download_button(
                "📄 Export Markdown",
                export_script_markdown(data),
                file_name=f"script_{datetime.now().strftime('%Y%m%d_%H%M')}.md",
                mime="text/markdown",
            )
        with exp_c2:
            st.download_button(
                "📝 Export Plain Text",
                export_script_plain(data),
                file_name=f"script_{datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                mime="text/plain",
            )
        with exp_c3:
            st.download_button(
                "🎬 Export SRT Subtitles",
                generate_srt(dialogue),
                file_name=f"subtitles_{datetime.now().strftime('%Y%m%d_%H%M')}.srt",
                mime="text/srt",
            )

        with st.expander("✏️ Edit Script", expanded=False):
            with st.form("edit_form"):
                has_caller = any(l["speaker"] == "Caller" for l in dialogue)
                speakers = ["Host 1", "Host 2"] + (["Caller"] if has_caller else [])

                new_dialogue = []
                for i, line in enumerate(dialogue):
                    c1, c2 = st.columns([1, 5])
                    sp = c1.selectbox(
                        "Speaker",
                        speakers,
                        index=(
                            speakers.index(line["speaker"])
                            if line["speaker"] in speakers
                            else 0
                        ),
                        key=f"sp_{i}",
                    )
                    tx = c2.text_area("Line", line["text"], height=80, key=f"tx_{i}")
                    new_dialogue.append({"speaker": sp, "text": tx})

                if st.form_submit_button("💾 Save Edits"):
                    st.session_state.script_data["dialogue"] = new_dialogue
                    st.success("Saved.")
                    st.rerun()

        # --- Rehearsal ---
        st.subheader("🎧 Live Rehearsal")
        idx = st.selectbox(
            "Preview a line",
            range(len(dialogue)),
            format_func=lambda i: (
                f"[{dialogue[i]['speaker']}] {dialogue[i]['text'][:80]}…"
            ),
        )
        if st.button("▶️ Play Line"):
            line = dialogue[idx]
            voice = m_voice if line["speaker"] == "Host 1" else f_voice
            speed = host1_speed if line["speaker"] == "Host 1" else host2_speed
            if line["speaker"] == "Caller":
                voice = c_voice
                speed = caller_speed

            if use_edge:
                with tempfile.NamedTemporaryFile(suffix=".mp3", delete=False) as tmp:
                    with st.spinner("Generating voice…"):
                        if _sync_edge_tts(line["text"], voice, tmp.name, speed):
                            st.audio(tmp.name)
                        else:
                            st.error("Edge TTS preview failed.")
            else:
                if not audio_client:
                    st.error("OpenAI key required for TTS.")
                else:
                    with tempfile.NamedTemporaryFile(suffix=".mp3", delete=False) as tmp:
                        if generate_tts(
                            audio_client, line["text"], voice, tmp.name,
                            tts_model=tts_model, speed=speed,
                        ):
                            st.audio(tmp.name)
                        else:
                            st.error("TTS preview failed.")


# === TAB 4 — PRODUCTION =====================================================
with tab4:
    st.header("🎚️ Final Production")

    if not st.session_state.script_data:
        st.info("Generate a script in the **Script & Rehearsal** tab first.")

    elif st.button("🚀 Produce Final Podcast", type="primary"):
        if not use_edge and not openai_key:
            st.error("OpenAI key required for TTS production. Or switch to Edge TTS (Free).")
            st.stop()

        progress = st.progress(0, text="Starting production…")
        status = st.empty()

        tts_client = OpenAI(api_key=openai_key) if (not use_edge and openai_key) else None

        with tempfile.TemporaryDirectory() as tmp_dir:
            tmp = Path(tmp_dir)
            script = st.session_state.script_data["dialogue"]

            args_list = [
                (
                    i, line, tts_client, m_voice, f_voice, c_voice, tmp,
                    tts_model, host1_speed, host2_speed, caller_speed, use_edge,
                )
                for i, line in enumerate(script)
            ]

            tts_errors: List[str] = []
            completed = 0

            with ThreadPoolExecutor(max_workers=MAX_TTS_WORKERS) as pool:
                futures = {
                    pool.submit(_voice_line_worker, a): a[0] for a in args_list
                }
                for future in as_completed(futures):
                    i, ok, err_msg = future.result()
                    completed += 1
                    progress.progress(
                        completed / len(script),
                        text=f"Voicing line {completed}/{len(script)}…",
                    )
                    if not ok:
                        tts_errors.append(f"Line {i}: {err_msg}")

            for err in tts_errors:
                st.warning(err)

            status.text("🎛️ Mixing final podcast…")
            out_path = mix_final_audio(
                tmp_dir, script, bg_source, selected_bg_url,
                uploaded_bg_file, music_ramp_up, uploaded_intro, uploaded_outro,
            )

            if out_path and out_path.exists():
                with open(out_path, "rb") as f:
                    audio_bytes = f.read()

                progress.progress(1.0, text="✅ Complete!")
                status.empty()
                st.audio(audio_bytes, format="audio/mp3")

                dl_c1, dl_c2 = st.columns(2)
                with dl_c1:
                    st.download_button(
                        "⬇️ Download Podcast",
                        audio_bytes,
                        file_name=f"podcast_{datetime.now().strftime('%Y%m%d_%H%M')}.mp3",
                        mime="audio/mp3",
                    )
                with dl_c2:
                    st.download_button(
                        "🎬 Download Subtitles",
                        generate_srt(script),
                        file_name=f"podcast_{datetime.now().strftime('%Y%m%d_%H%M')}.srt",
                        mime="text/srt",
                    )

                duration_est = sum(len(l["text"].split()) for l in script) / 150
                size_mb = len(audio_bytes) / 1_048_576
                tts_label = "Edge TTS (free)" if use_edge else ("HD audio" if tts_hd else "Standard TTS")
                st.success(
                    f"🎉 Done! ~{duration_est:.0f} min · {size_mb:.1f} MB · {tts_label}"
                )
            else:
                st.error("Production failed — check warnings above.")


# === TAB 5 — STUDY TOOLS ====================================================
with tab5:
    st.header("📚 Study Tools")
    st.caption(
        "Generate study materials from your source content — "
        "similar to Google NotebookLM, but open source and using your own LLM."
    )

    if not st.session_state.source_text:
        st.info("Load source content in the **Source** tab first.")
    else:
        # Count how many tools have been generated
        tool_state_keys = {
            "Study Guide": "study_guide",
            "Briefing Document": "briefing_doc",
            "FAQ": "faq_doc",
            "Flashcards": "flashcards",
            "Timeline": "timeline",
            "Key Concepts": "key_concepts",
        }
        generated_count = sum(
            1 for k in tool_state_keys.values() if st.session_state.get(k)
        )
        st.caption(f"{generated_count}/{len(tool_state_keys)} tools generated for current source.")

        # --- Generate All button ---
        if st.button("🚀 Generate All Study Tools", type="primary"):
            client, model, err = get_llm_client(
                model_choice, _resolve_model_name(),
                openai_key, xai_key, deepseek_key, budget_mode,
            )
            if err:
                st.error(err)
            else:
                progress = st.progress(0, text="Generating study tools…")
                tools_list = list(tool_state_keys.items())

                for idx, (tool_name, state_key) in enumerate(tools_list):
                    if st.session_state.get(state_key):
                        progress.progress(
                            (idx + 1) / len(tools_list),
                            text=f"Skipping {tool_name} (already generated)…",
                        )
                        continue

                    progress.progress(
                        (idx + 1) / len(tools_list),
                        text=f"Generating {tool_name}…",
                    )

                    is_json = tool_name == "Flashcards"
                    result = generate_study_tool(
                        client, model, st.session_state.source_text,
                        tool_name, language, is_json=is_json,
                    )

                    if result:
                        if is_json:
                            try:
                                cards = json.loads(result)
                                if isinstance(cards, dict) and "flashcards" in cards:
                                    cards = cards["flashcards"]
                                st.session_state[state_key] = cards
                            except json.JSONDecodeError:
                                st.session_state[state_key] = result
                        else:
                            st.session_state[state_key] = result

                progress.progress(1.0, text="✅ All tools generated!")
                st.rerun()

        st.divider()

        # --- Individual tools in expanders ---

        # 1. Study Guide
        st.subheader("📖 Study Guide")
        col_gen, col_dl = st.columns([3, 1])
        with col_gen:
            if st.button("Generate Study Guide", key="gen_study_guide"):
                client, model, err = get_llm_client(
                    model_choice, _resolve_model_name(),
                    openai_key, xai_key, deepseek_key, budget_mode,
                )
                if err:
                    st.error(err)
                else:
                    with st.spinner("Generating study guide…"):
                        result = generate_study_tool(
                            client, model, st.session_state.source_text,
                            "Study Guide", language,
                        )
                        if result:
                            st.session_state.study_guide = result
                            st.rerun()

        if st.session_state.study_guide:
            with col_dl:
                st.download_button(
                    "⬇️ Download",
                    st.session_state.study_guide,
                    file_name=f"study_guide_{datetime.now().strftime('%Y%m%d_%H%M')}.md",
                    mime="text/markdown",
                    key="dl_study_guide",
                )
            with st.expander("View Study Guide", expanded=True):
                st.markdown(st.session_state.study_guide)

        st.divider()

        # 2. Briefing Document
        st.subheader("📋 Briefing Document")
        col_gen, col_dl = st.columns([3, 1])
        with col_gen:
            if st.button("Generate Briefing Document", key="gen_briefing"):
                client, model, err = get_llm_client(
                    model_choice, _resolve_model_name(),
                    openai_key, xai_key, deepseek_key, budget_mode,
                )
                if err:
                    st.error(err)
                else:
                    with st.spinner("Generating briefing document…"):
                        result = generate_study_tool(
                            client, model, st.session_state.source_text,
                            "Briefing Document", language,
                        )
                        if result:
                            st.session_state.briefing_doc = result
                            st.rerun()

        if st.session_state.briefing_doc:
            with col_dl:
                st.download_button(
                    "⬇️ Download",
                    st.session_state.briefing_doc,
                    file_name=f"briefing_{datetime.now().strftime('%Y%m%d_%H%M')}.md",
                    mime="text/markdown",
                    key="dl_briefing",
                )
            with st.expander("View Briefing Document", expanded=True):
                st.markdown(st.session_state.briefing_doc)

        st.divider()

        # 3. FAQ
        st.subheader("❓ FAQ")
        col_gen, col_dl = st.columns([3, 1])
        with col_gen:
            if st.button("Generate FAQ", key="gen_faq"):
                client, model, err = get_llm_client(
                    model_choice, _resolve_model_name(),
                    openai_key, xai_key, deepseek_key, budget_mode,
                )
                if err:
                    st.error(err)
                else:
                    with st.spinner("Generating FAQ…"):
                        result = generate_study_tool(
                            client, model, st.session_state.source_text,
                            "FAQ", language,
                        )
                        if result:
                            st.session_state.faq_doc = result
                            st.rerun()

        if st.session_state.faq_doc:
            with col_dl:
                st.download_button(
                    "⬇️ Download",
                    st.session_state.faq_doc,
                    file_name=f"faq_{datetime.now().strftime('%Y%m%d_%H%M')}.md",
                    mime="text/markdown",
                    key="dl_faq",
                )
            with st.expander("View FAQ", expanded=True):
                st.markdown(st.session_state.faq_doc)

        st.divider()

        # 4. Flashcards
        st.subheader("🃏 Flashcards")
        col_gen, col_dl = st.columns([3, 1])
        with col_gen:
            if st.button("Generate Flashcards", key="gen_flashcards"):
                client, model, err = get_llm_client(
                    model_choice, _resolve_model_name(),
                    openai_key, xai_key, deepseek_key, budget_mode,
                )
                if err:
                    st.error(err)
                else:
                    with st.spinner("Generating flashcards…"):
                        result = generate_study_tool(
                            client, model, st.session_state.source_text,
                            "Flashcards", language, is_json=True,
                        )
                        if result:
                            try:
                                cards = json.loads(result)
                                if isinstance(cards, dict) and "flashcards" in cards:
                                    cards = cards["flashcards"]
                                st.session_state.flashcards = cards
                            except json.JSONDecodeError:
                                st.session_state.flashcards = result
                            st.rerun()

        if st.session_state.flashcards:
            cards = st.session_state.flashcards
            if isinstance(cards, list):
                with col_dl:
                    st.download_button(
                        "⬇️ Download",
                        render_flashcards_markdown(cards),
                        file_name=f"flashcards_{datetime.now().strftime('%Y%m%d_%H%M')}.md",
                        mime="text/markdown",
                        key="dl_flashcards",
                    )
                with st.expander(f"View Flashcards ({len(cards)} cards)", expanded=True):
                    # Interactive flashcard viewer
                    card_idx = st.selectbox(
                        "Select card",
                        range(len(cards)),
                        format_func=lambda i: f"Card {i + 1}: {str(cards[i].get('front', cards[i].get('Front', '?')))[:60]}…",
                        key="flashcard_select",
                    )
                    card = cards[card_idx]
                    front = card.get("front") or card.get("Front") or card.get("question") or "?"
                    back = card.get("back") or card.get("Back") or card.get("answer") or "?"

                    st.markdown(f"### 📌 {front}")
                    if st.button("🔄 Reveal Answer", key="reveal_flashcard"):
                        st.markdown(f"**Answer:** {back}")
                    else:
                        st.caption("Click 'Reveal Answer' to see the back of the card.")
            else:
                # Fallback if flashcards came back as text instead of JSON
                with st.expander("View Flashcards", expanded=True):
                    st.markdown(str(cards))

        st.divider()

        # 5. Timeline
        st.subheader("📅 Timeline")
        col_gen, col_dl = st.columns([3, 1])
        with col_gen:
            if st.button("Generate Timeline", key="gen_timeline"):
                client, model, err = get_llm_client(
                    model_choice, _resolve_model_name(),
                    openai_key, xai_key, deepseek_key, budget_mode,
                )
                if err:
                    st.error(err)
                else:
                    with st.spinner("Generating timeline…"):
                        result = generate_study_tool(
                            client, model, st.session_state.source_text,
                            "Timeline", language,
                        )
                        if result:
                            st.session_state.timeline = result
                            st.rerun()

        if st.session_state.timeline:
            with col_dl:
                st.download_button(
                    "⬇️ Download",
                    st.session_state.timeline,
                    file_name=f"timeline_{datetime.now().strftime('%Y%m%d_%H%M')}.md",
                    mime="text/markdown",
                    key="dl_timeline",
                )
            with st.expander("View Timeline", expanded=True):
                st.markdown(st.session_state.timeline)

        st.divider()

        # 6. Key Concepts / Glossary
        st.subheader("📘 Key Concepts & Glossary")
        col_gen, col_dl = st.columns([3, 1])
        with col_gen:
            if st.button("Generate Key Concepts", key="gen_concepts"):
                client, model, err = get_llm_client(
                    model_choice, _resolve_model_name(),
                    openai_key, xai_key, deepseek_key, budget_mode,
                )
                if err:
                    st.error(err)
                else:
                    with st.spinner("Generating key concepts…"):
                        result = generate_study_tool(
                            client, model, st.session_state.source_text,
                            "Key Concepts", language,
                        )
                        if result:
                            st.session_state.key_concepts = result
                            st.rerun()

        if st.session_state.key_concepts:
            with col_dl:
                st.download_button(
                    "⬇️ Download",
                    st.session_state.key_concepts,
                    file_name=f"glossary_{datetime.now().strftime('%Y%m%d_%H%M')}.md",
                    mime="text/markdown",
                    key="dl_concepts",
                )
            with st.expander("View Key Concepts", expanded=True):
                st.markdown(st.session_state.key_concepts)

        # --- Export all tools as a single document ---
        st.divider()
        if generated_count > 0:
            all_tools_parts = [
                f"# Study Materials\n",
                f"**Generated:** {datetime.now().strftime('%Y-%m-%d %H:%M')}\n",
                f"**Source:** {len(st.session_state.source_text):,} characters\n\n",
            ]
            if st.session_state.study_guide:
                all_tools_parts.append("---\n\n# Study Guide\n\n")
                all_tools_parts.append(st.session_state.study_guide + "\n\n")
            if st.session_state.briefing_doc:
                all_tools_parts.append("---\n\n# Briefing Document\n\n")
                all_tools_parts.append(st.session_state.briefing_doc + "\n\n")
            if st.session_state.faq_doc:
                all_tools_parts.append("---\n\n# FAQ\n\n")
                all_tools_parts.append(st.session_state.faq_doc + "\n\n")
            if st.session_state.flashcards and isinstance(st.session_state.flashcards, list):
                all_tools_parts.append("---\n\n")
                all_tools_parts.append(render_flashcards_markdown(st.session_state.flashcards) + "\n\n")
            if st.session_state.timeline:
                all_tools_parts.append("---\n\n# Timeline\n\n")
                all_tools_parts.append(st.session_state.timeline + "\n\n")
            if st.session_state.key_concepts:
                all_tools_parts.append("---\n\n# Key Concepts & Glossary\n\n")
                all_tools_parts.append(st.session_state.key_concepts + "\n\n")

            combined = "\n".join(all_tools_parts)
            st.download_button(
                "📦 Download All Study Materials",
                combined,
                file_name=f"study_materials_{datetime.now().strftime('%Y%m%d_%H%M')}.md",
                mime="text/markdown",
                type="primary",
            )
